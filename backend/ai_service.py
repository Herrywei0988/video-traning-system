"""
ai_service.py — 支援多種檔案類型的 AI 分析服務
- 影音：MP4, MOV, AVI, MKV, WebM, MP3, WAV, M4A ...
- 文件：TXT, PDF, DOCX
- 投影片：PPTX
"""
import os
import json
import subprocess
import re
from pathlib import Path
from openai import OpenAI

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
MODEL = os.getenv("OPENAI_MODEL", "gpt-4o")

VIDEO_EXTS  = {'.mp4', '.mov', '.avi', '.mkv', '.webm', '.flv', '.wmv', '.m4v'}
AUDIO_EXTS  = {'.mp3', '.wav', '.m4a', '.aac', '.ogg', '.flac', '.opus', '.wma'}
TEXT_EXTS   = {'.txt', '.md', '.csv'}
PDF_EXTS    = {'.pdf'}
DOCX_EXTS   = {'.docx', '.doc'}
PPTX_EXTS   = {'.pptx', '.ppt'}
ALL_SUPPORTED = VIDEO_EXTS | AUDIO_EXTS | TEXT_EXTS | PDF_EXTS | DOCX_EXTS | PPTX_EXTS

def classify_file(filepath: str) -> str:
    suffix = Path(filepath).suffix.lower()
    if suffix in VIDEO_EXTS:  return 'video'
    if suffix in AUDIO_EXTS:  return 'audio'
    if suffix in PPTX_EXTS:   return 'pptx'
    return 'document'

def extract_text_from_txt(filepath: str):
    for enc in ['utf-8', 'utf-8-sig', 'big5', 'gbk', 'latin-1']:
        try:
            with open(filepath, 'r', encoding=enc) as f:
                text = f.read()
            return text, text.count('\n') + 1
        except Exception:
            continue
    raise ValueError("無法讀取文字檔（不支援的編碼）")

def extract_text_from_pdf(filepath: str):
    import pdfplumber
    pages = []
    with pdfplumber.open(filepath) as pdf:
        page_count = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[第 {i+1} 頁]\n{text.strip()}")
    return "\n\n".join(pages), page_count

def extract_text_from_docx(filepath: str):
    from docx import Document
    doc = Document(filepath)
    sections = []
    para_count = 0
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        para_count += 1
        style = para.style.name if para.style else ''
        if 'Heading' in style or '標題' in style:
            sections.append(f"\n## {text}")
        else:
            sections.append(text)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sections.append("\n[表格]\n" + "\n".join(rows))
    return "\n".join(sections), para_count

def extract_text_from_pptx(filepath: str):
    from pptx import Presentation
    prs = Presentation(filepath)
    slides_text = []
    slide_count = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        parts = []
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f"### 投影片 {i+1}：{slide.shapes.title.text.strip()}")
        else:
            parts.append(f"### 投影片 {i+1}")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape == slide.shapes.title:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    parts.append(text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[備註] {notes}")
        if len(parts) > 1:
            slides_text.append("\n".join(parts))
    return "\n\n".join(slides_text), slide_count

def extract_audio(video_path: str) -> str:
    suffix = Path(video_path).suffix.lower()
    if suffix in AUDIO_EXTS:
        return video_path
    audio_path = video_path.rsplit('.', 1)[0] + '_audio.mp3'
    try:
        result = subprocess.run([
            'ffmpeg', '-i', video_path, '-vn', '-acodec', 'mp3', '-ab', '128k', '-y', audio_path
        ], capture_output=True, text=True, timeout=300)
        if result.returncode != 0:
            raise RuntimeError(f"ffmpeg 轉檔失敗：{result.stderr[:300]}")
        return audio_path
    except FileNotFoundError:
        raise RuntimeError(
            "系統未安裝 ffmpeg，無法處理影片檔案。"
            "請在伺服器執行：sudo apt install ffmpeg（Linux/WSL）"
            "或 brew install ffmpeg（macOS）"
        )

def get_audio_duration(path: str) -> float:
    try:
        result = subprocess.run([
            'ffprobe', '-v', 'quiet', '-print_format', 'json', '-show_streams', path
        ], capture_output=True, text=True, timeout=30)
        data = json.loads(result.stdout)
        for stream in data.get('streams', []):
            if 'duration' in stream:
                return float(stream['duration'])
    except Exception:
        pass
    return 0.0

def chunk_audio(audio_path: str, max_size_mb: int = 20) -> list:
    """Returns: list of (chunk_path, start_offset_seconds)"""
    file_size_mb = os.path.getsize(audio_path) / (1024 * 1024)
    if file_size_mb <= max_size_mb:
        return [(audio_path, 0.0)]
    duration = get_audio_duration(audio_path)
    if duration == 0:
        return [(audio_path, 0.0)]
    chunk_duration = max(30, int((max_size_mb / file_size_mb) * duration * 0.9))
    chunks, start, idx = [], 0, 0
    base = audio_path.rsplit('.', 1)[0]
    while start < duration:
        chunk_path = f"{base}_chunk{idx}.mp3"
        subprocess.run([
            'ffmpeg', '-i', audio_path, '-ss', str(start),
            '-t', str(chunk_duration), '-acodec', 'mp3', '-y', chunk_path
        ], capture_output=True, timeout=120)
        if os.path.exists(chunk_path):
            chunks.append((chunk_path, float(start)))
        start += chunk_duration
        idx += 1
    return chunks if chunks else [(audio_path, 0.0)]

def transcribe_audio(filepath: str) -> tuple:
    """Returns: (full_text, segments_list)
    segments_list: [{"start": float, "end": float, "text": str}, ...]
    """
    audio_path = extract_audio(filepath)
    chunks = chunk_audio(audio_path)
    
    full_text_parts = []
    all_segments = []
    
    for chunk_path, start_offset in chunks:
        with open(chunk_path, 'rb') as f:
            response = client.audio.transcriptions.create(
                model="whisper-1",
                file=f,
                language="zh",
                response_format="verbose_json",
                timestamp_granularities=["segment"]
            )
        
        full_text_parts.append(response.text)
        
        for seg in response.segments:
            all_segments.append({
                "start": round(seg.start + start_offset, 2),
                "end": round(seg.end + start_offset, 2),
                "text": seg.text.strip()
            })
        
        # Cleanup temp chunks
        if chunk_path != filepath and chunk_path != audio_path and '_chunk' in chunk_path:
            try: os.remove(chunk_path)
            except Exception: pass
    
    if audio_path != filepath and os.path.exists(audio_path):
        try: os.remove(audio_path)
        except Exception: pass
    
    return "\n".join(full_text_parts), all_segments

REFUSAL_INDICATORS = [
    "i'm sorry", "i am sorry", "i cannot", "i can't", "i can not",
    "unable to", "not able to", "對不起", "抱歉", "無法協助", "無法處理",
]

def _is_refusal(text: str) -> bool:
    snippet = text[:200].lower()
    # 如果前 100 字就有 { 代表是 JSON 開頭，不是拒絕
    if '{' in text[:100]:
        return False
    return any(ind in snippet for ind in REFUSAL_INDICATORS)

def _build_prompt(file_type: str, title: str, description: str, content: str, has_timestamps: bool = False) -> str:
    if file_type == 'pptx':
        label = "投影片內容（逐頁整理）"
        hint = "請特別注意每張投影片的主題與邏輯脈絡，整理出整份簡報的論述結構與每頁重點。"
    elif file_type == 'document':
        label = "文件內容"
        hint = "請根據文件段落結構整理重點，若有表格請說明核心資訊。"
    else:
        label = "逐字稿（每行前方 [MM:SS] 為該段開始時間）" if has_timestamps else "逐字稿"
        hint = "請根據口語脈絡整理重點，注意講者強調的部分。"

    if has_timestamps:
        segments_schema = '{"title": "段落主題", "content": "段落摘要（50-80字）", "importance": "高/中/低", "start_time": 秒數}'
        segments_rule = """
特別注意：key_segments 的每一筆必須包含 "start_time" 欄位，格式為秒數（整數或小數，例如 223 代表第 3 分 43 秒）。請從上方逐字稿的 [MM:SS] 時間戳中，找出最能代表該段落重點的時間點，轉成秒數填入。這是用來讓使用者點擊段落直接跳到影片對應位置的。"""
    else:
        segments_schema = '{"title": "段落主題", "content": "段落摘要（50-80字）", "importance": "高/中/低"}'
        segments_rule = ""

    return f"""你是一位專業的教育培訓內容整理專家，負責幫助補教機構把培訓資料整理成可執行的知識模組。

請根據以下{label}，產出完整的結構化分析。請用繁體中文，嚴格按照以下 JSON 格式輸出，不加任何額外說明或 markdown：

{{
  "summary": "整體摘要（200-300字，說明本次培訓的核心目標與主要內容）",
  "key_points": [
    {{"point": "重點標題", "detail": "詳細說明"}},
    ...
  ],
  "action_items": [
    {{"task": "具體待辦事項", "owner": "建議負責人（例如班主任/老師/主管）", "priority": "高/中/低"}},
    ...
  ],
  "faq": [
    {{"question": "常見問題", "answer": "建議回答方式"}},
    ...
  ],
  "key_segments": [
    {segments_schema},
    ...
  ],
  "exec_summary": "主管版摘要（100字以內，聚焦在需要決策的事項、管理重點、整體方向）",
  "manager_summary": "班主任版摘要（150字以內，聚焦在落地執行方法、日常管理要點、需要帶領老師做的事）",
  "teacher_summary": "老師版重點（150字以內，聚焦在課堂執行、與學生互動、具體教學動作）",
  "topics": ["從以下選擇相關主題：招生, 續約, 教學, 班務, 家長溝通, 培訓, 行政, 品質管理, 其他"]
}}
{segments_rule}

{hint}
資料標題：{title}
資料說明：{description or '（無說明）'}

{label}：
{content}

請確保每個段落都有實質內容，避免空泛描述。key_points 至少 3 條，action_items 至少 3 條，key_segments 至少 4 條。"""

def analyze(content: str, title: str, description: str, file_type: str = 'video', segments: list = None) -> dict:
    has_timestamps = bool(segments) and file_type in ('video', 'audio')
    
    if has_timestamps:
        lines = []
        for seg in segments:
            mm = int(seg["start"] // 60)
            ss = int(seg["start"] % 60)
            lines.append(f"[{mm:02d}:{ss:02d}] {seg['text']}")
        prompt_content = "\n".join(lines)
    else:
        prompt_content = content
    
    # gpt-4o 有 128k context，可以放更多內容
    truncated = prompt_content[:15000] if len(prompt_content) > 15000 else prompt_content
    
    prompt = _build_prompt(file_type, title, description, truncated, has_timestamps=has_timestamps)
    
    response = client.chat.completions.create(
        model=MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_tokens=3000
    )
    raw = response.choices[0].message.content.strip()

    # Detect AI refusal before trying to parse JSON
    if _is_refusal(raw):
        raise ValueError(
            "AI 判定此內容不適合做為培訓資料整理。"
            "可能原因：內容過於簡短、非培訓主題的閒聊內容、或資訊密度不足。"
            "建議提供包含明確培訓目標、操作步驟或具體知識點的資料。"
        )

    if raw.startswith("```"):
        lines_raw = raw.split('\n')
        raw = '\n'.join(lines_raw[1:-1]) if lines_raw[-1] == '```' else '\n'.join(lines_raw[1:])
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        match = re.search(r'\{.*\}', raw, re.DOTALL)
        if match:
            return json.loads(match.group())
        raise ValueError(f"Could not parse AI response: {raw[:200]}")

def process_file(filepath: str, title: str, description: str) -> tuple:
    """Returns: (file_type, content_text, analysis_dict, duration_or_none, page_count_or_none, segments_or_none)"""
    suffix = Path(filepath).suffix.lower()
    ft = classify_file(filepath)

    if suffix in VIDEO_EXTS or suffix in AUDIO_EXTS:
        duration = get_audio_duration(filepath)
        transcript, segments = transcribe_audio(filepath)
        analysis = analyze(transcript, title, description, file_type=ft, segments=segments)
        return ft, transcript, analysis, duration, None, segments

    if suffix in TEXT_EXTS:
        text, lines = extract_text_from_txt(filepath)
        analysis = analyze(text, title, description, file_type='document')
        return 'document', text, analysis, None, lines, None

    if suffix in PDF_EXTS:
        text, pages = extract_text_from_pdf(filepath)
        if not text.strip():
            raise ValueError("PDF 中找不到可讀取的文字（可能是掃描檔，目前不支援 OCR）")
        analysis = analyze(text, title, description, file_type='document')
        return 'document', text, analysis, None, pages, None

    if suffix in DOCX_EXTS:
        text, paras = extract_text_from_docx(filepath)
        analysis = analyze(text, title, description, file_type='document')
        return 'document', text, analysis, None, paras, None

    if suffix in PPTX_EXTS:
        text, slides = extract_text_from_pptx(filepath)
        analysis = analyze(text, title, description, file_type='pptx')
        return 'pptx', text, analysis, None, slides, None

    raise ValueError(f"不支援的檔案格式：{suffix}")